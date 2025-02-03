import requests
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
import graphviz
import matplotlib.pyplot as plt
import numpy as np
import seaborn as sns

def download_image(url, filename, images_dir):
    """Download image from URL and save to images directory"""
    try:
        response = requests.get(url, stream=True)
        response.raise_for_status()
        
        file_path = images_dir / filename
        with open(file_path, 'wb') as f:
            for chunk in response.iter_content(chunk_size=8192):
                f.write(chunk)
        return file_path
    except Exception as e:
        print(f"Error downloading {filename}: {e}")
        return None

def setup_image_directory():
    """Create images directory if it doesn't exist"""
    images_dir = Path('presentation_images')
    images_dir.mkdir(exist_ok=True)
    return images_dir

def download_presentation_images():
    """Download all images needed for the presentation"""
    images_dir = setup_image_directory()
    
    # Dictionary of image URLs and their local filenames
    images = {
        'random_forest.jpg': 'https://raw.githubusercontent.com/scikit-learn/scikit-learn/main/doc/images/plot_randomized_forest.png',
        'pipeline.png': 'https://scikit-learn.org/stable/_images/grid_search_workflow.png',
        'preprocessing.png': 'https://raw.githubusercontent.com/scikit-learn/scikit-learn/main/doc/images/plot_preprocessing.png',
        'concept_drift.png': 'https://raw.githubusercontent.com/scikit-learn/scikit-learn/main/doc/images/plot_incremental_learning.png',
        'evaluation.png': 'https://scikit-learn.org/stable/_images/sphx_glr_plot_confusion_matrix_001.png'
    }
    
    downloaded_images = {}
    for filename, url in images.items():
        file_path = download_image(url, filename, images_dir)
        if file_path:
            downloaded_images[filename] = file_path
            print(f"Successfully downloaded {filename}")
    
    return downloaded_images

def add_title_slide(prs, title, subtitle):
    """Add a title slide with subtitle"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    # Format title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True

def add_content_slide_with_image(prs, title, content, image_path=None, layout_idx=1):
    """Add a content slide with bullet points and optional image"""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    
    # Add and format title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    
    # Add content
    if image_path:
        # Adjust content width to make room for image
        left = Inches(0.5)
        width = Inches(4)
        height = Inches(5.5)
        
        # Add image
        img_left = Inches(5)
        img_top = Inches(1.5)
        img_width = Inches(4)
        
        if image_path.exists():
            slide.shapes.add_picture(str(image_path), img_left, img_top, width=img_width)
    else:
        left = Inches(1)
        width = Inches(8)
        height = Inches(5.5)
    
    # Add text box for content
    txBox = slide.shapes.add_textbox(left, Inches(1.5), width, height)
    tf = txBox.text_frame
    
    for point in content:
        p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(24)
        p.level = 0

def create_pipeline_diagram():
    """Create a custom pipeline diagram using graphviz"""
    dot = graphviz.Digraph(comment='Model Pipeline')
    dot.attr(rankdir='TB')  # Top to Bottom direction
    dot.attr('node', shape='rectangle', style='rounded,filled', fillcolor='lightgray')
    
    # Text Preprocessor Cluster
    with dot.subgraph(name='cluster_0') as c:
        c.attr(label='1. Text Preprocessor', style='rounded', color='blue', fontcolor='blue')
        c.attr('node', fillcolor='lightblue')
        
        # Create nodes for each preprocessing step
        c.node('url', 'URL & Code Block\nRemoval')
        c.node('norm', 'Text Normalization\n(lowercase, whitespace)')
        c.node('token', 'Tokenization')
        c.node('stop', 'Stop-word Removal')
        c.node('lemma', 'Lemmatization')
        
        # Connect preprocessing steps
        c.edge('url', 'norm')
        c.edge('norm', 'token')
        c.edge('token', 'stop')
        c.edge('stop', 'lemma')
    
    # TF-IDF Vectorizer Cluster
    with dot.subgraph(name='cluster_1') as c:
        c.attr(label='2. TF-IDF Vectorizer', style='rounded', color='darkgreen', fontcolor='darkgreen')
        c.attr('node', fillcolor='lightgreen')
        
        c.node('tfidf_params', 'Parameters:\nmax_features=5000')
        c.node('tfidf_process', 'Convert Text to\nNumerical Features')
        
        c.edge('tfidf_params', 'tfidf_process')
    
    # Random Forest Cluster
    with dot.subgraph(name='cluster_2') as c:
        c.attr(label='3. Random Forest Classifier', style='rounded', color='darkred', fontcolor='darkred')
        c.attr('node', fillcolor='lightpink')
        
        c.node('rf_params', 'Configuration:\nn_estimators=100\nmax_depth=None\nrandom_state=42')
        c.node('rf_train', 'Train Model')
        c.node('rf_predict', 'Make Predictions')
        
        c.edge('rf_params', 'rf_train')
        c.edge('rf_train', 'rf_predict')
    
    # Connect main components
    dot.edge('lemma', 'tfidf_process')
    dot.edge('tfidf_process', 'rf_train')
    
    # Add benefit notes
    dot.node('benefit1', 'Benefit: Clean and\nconsistent input', shape='note', fillcolor='lightyellow')
    dot.node('benefit2', 'Benefit: Highlights important\nwords, reduces noise', shape='note', fillcolor='lightyellow')
    dot.node('benefit3', 'Benefit: Handles large feature\nspaces robustly', shape='note', fillcolor='lightyellow')
    
    # Connect benefits to their components
    dot.edge('lemma', 'benefit1', style='dashed')
    dot.edge('tfidf_process', 'benefit2', style='dashed')
    dot.edge('rf_predict', 'benefit3', style='dashed')
    
    # Save diagram with higher DPI for better quality
    dot.attr(dpi='300')
    dot.render('presentation_images/custom_pipeline', format='png', cleanup=True)
    return Path('presentation_images/custom_pipeline.png')

def create_training_flow_diagram():
    """Create a diagram showing the training process flow"""
    dot = graphviz.Digraph(comment='Training Flow')
    dot.attr(rankdir='TB')  # Top to Bottom direction
    
    # Data Loading & Preparation
    with dot.subgraph(name='cluster_0') as c:
        c.attr(label='Data Preparation')
        c.node('D1', 'Load GitHub\nIssues Dataset')
        c.node('D2', 'Clean & Preprocess\nText Data')
        c.node('D3', 'TF-IDF\nVectorization')
        
        c.edge('D1', 'D2')
        c.edge('D2', 'D3')
    
    # Model Training
    with dot.subgraph(name='cluster_1') as c:
        c.attr(label='Model Training')
        c.node('T1', 'Train-Test Split\n(80-20)')
        c.node('T2', 'Random Forest\nTraining')
        c.node('T3', 'Cross-Validation')
        
        c.edge('T1', 'T2')
        c.edge('T2', 'T3')
    
    # Evaluation
    with dot.subgraph(name='cluster_2') as c:
        c.attr(label='Evaluation')
        c.node('E1', 'Performance\nMetrics')
        c.node('E2', 'Model\nSerialization')
        
        c.edge('E1', 'E2')
    
    # Connect components
    dot.edge('D3', 'T1')
    dot.edge('T3', 'E1')
    
    # Save diagram
    dot.render('presentation_images/training_flow', format='png', cleanup=True)
    return Path('presentation_images/training_flow.png')

def create_concept_drift_diagram():
    """Create a visualization of concept drift impact"""
    plt.figure(figsize=(10, 6))
    
    # Generate sample data
    x = np.linspace(0, 10, 100)
    y1 = np.sin(x) + np.random.normal(0, 0.1, 100)
    y2 = np.sin(x + 1) + np.random.normal(0, 0.1, 100)
    y3 = np.sin(x + 2) + np.random.normal(0, 0.1, 100)
    
    plt.plot(x, y1, label='Initial Data Distribution')
    plt.plot(x, y2, label='Distribution Shift 1')
    plt.plot(x, y3, label='Distribution Shift 2')
    
    plt.title('Concept Drift in GitHub Issues Over Time')
    plt.xlabel('Time')
    plt.ylabel('Data Distribution')
    plt.legend()
    plt.grid(True)
    
    # Save plot
    plt.savefig('presentation_images/concept_drift_custom.png')
    plt.close()
    return Path('presentation_images/concept_drift_custom.png')

def create_performance_visualization():
    """Create custom performance visualization"""
    # Sample performance data
    categories = ['Bugs', 'Enhancements', 'Questions']
    precision = [0.76, 0.70, 0.61]
    recall = [0.79, 0.80, 0.09]
    f1_scores = [0.77, 0.75, 0.16]
    
    # Set up the plot
    plt.figure(figsize=(10, 6))
    x = np.arange(len(categories))
    width = 0.25
    
    plt.bar(x - width, precision, width, label='Precision')
    plt.bar(x, recall, width, label='Recall')
    plt.bar(x + width, f1_scores, width, label='F1-Score')
    
    plt.ylabel('Scores')
    plt.title('Model Performance by Issue Type')
    plt.xticks(x, categories)
    plt.legend()
    plt.grid(True, alpha=0.3)
    
    # Save plot
    plt.savefig('presentation_images/performance_custom.png')
    plt.close()
    return Path('presentation_images/performance_custom.png')

def create_evaluation_diagram():
    """Create a compact evaluation diagram"""
    dot = graphviz.Digraph(comment='Model Evaluation')
    dot.attr(rankdir='LR')  # Left to Right direction
    dot.attr('node', shape='rectangle', style='rounded,filled')
    
    # Main process
    with dot.subgraph(name='cluster_0') as c:
        c.attr(label='Random Forest Training & Evaluation', style='rounded', color='navy', fontcolor='navy')
        
        # Data Split Cluster
        with dot.subgraph(name='cluster_split') as d:
            d.attr(label='Data Preparation & Split', style='rounded', color='blue', fontcolor='blue')
            d.attr('node', fillcolor='lightblue')
            
            d.node('dataset', '''GitHub Issues Dataset
-----------------
X: issue_body (text)
y: issue_label''')
            
            d.node('prep', '''Data Preparation
-----------------
• Handle missing values
X = df['issue_body'].fillna('')
y = df['issue_label']''')
            
            d.node('split', '''train_test_split
-----------------
test_size=0.2
random_state=42''')
            
            d.node('train', '''Training Data (80%)
-----------------
X_train, y_train''')
            
            d.node('test', '''Test Data (20%)
-----------------
X_test, y_test''')
            
            # Connect split nodes
            d.edge('dataset', 'prep')
            d.edge('prep', 'split')
            d.edge('split', 'train')
            d.edge('split', 'test')
        
        # Model Training Cluster
        with dot.subgraph(name='cluster_train') as t:
            t.attr(label='Model Training', style='rounded', color='darkgreen', fontcolor='darkgreen')
            t.attr('node', fillcolor='lightgreen')
            
            t.node('rf_config', '''Random Forest\nConfiguration
-----------------
• n_estimators=100
• max_depth=None
• random_state=42''')
            
            t.node('cv', '''Cross-Validation
-----------------
• K-fold validation
• Performance check''')
        
        # Evaluation & Deployment Cluster
        with dot.subgraph(name='cluster_eval') as e:
            e.attr(label='Evaluation & Deployment', style='rounded', color='darkred', fontcolor='darkred')
            e.attr('node', fillcolor='lightpink')
            
            e.node('eval', '''Model\nEvaluation''')
            
            e.node('metrics', '''Performance Metrics
-------------------
Accuracy: 73%
Bug (P/R): 76%/79%
Enhance (P/R): 70%/80%
Question (P/R): 61%/9%''',
            shape='note', fillcolor='lightyellow')
            
            e.node('serialize', '''Model Serialization
-----------------
joblib.dump(pipeline)''')
        
        # Training Methodology Note
        dot.node('methodology', '''Code Implementation
-----------------
X_train, X_test, y_train, y_test = 
    train_test_split(
        X, y,
        test_size=0.2,
        random_state=42
    )''',
        shape='note', fillcolor='lightyellow')
        
        # Connect all components
        dot.edge('train', 'rf_config')
        dot.edge('rf_config', 'cv')
        dot.edge('cv', 'eval')
        dot.edge('test', 'eval')
        dot.edge('eval', 'metrics')
        dot.edge('eval', 'serialize')
        dot.edge('split', 'methodology', style='dashed')
    
    # Save diagram with higher DPI and landscape orientation
    dot.attr(dpi='300')
    dot.attr(size='12,6')
    dot.render('presentation_images/evaluation_diagram', format='png', cleanup=True)
    return Path('presentation_images/evaluation_diagram.png')

def create_sgd_diagram():
    """Create a basic diagram showing SGD classifier concept"""
    dot = graphviz.Digraph(comment='SGD Classifier Basic')
    dot.attr(rankdir='LR')  # Left to Right direction
    dot.attr('node', shape='rectangle', style='rounded,filled')
    
    # Main process
    with dot.subgraph(name='cluster_0') as c:
        c.attr(label='SGD Classifier vs Random Forest', style='rounded', color='navy', fontcolor='navy')
        
        # Random Forest side
        with dot.subgraph(name='cluster_rf') as rf:
            rf.attr(label='Random Forest', style='rounded', color='darkred', fontcolor='darkred')
            rf.attr('node', fillcolor='lightpink')
            
            rf.node('rf_train', '''Traditional Training
-----------------
• Full dataset required
• Complete retraining
• Static model''')
            
            rf.node('rf_limit', '''Limitations
-----------------
• High memory usage
• Slow updates
• Fixed patterns''')
        
        # SGD side
        with dot.subgraph(name='cluster_sgd') as sgd:
            sgd.attr(label='SGD Classifier', style='rounded', color='darkgreen', fontcolor='darkgreen')
            sgd.attr('node', fillcolor='lightgreen')
            
            sgd.node('sgd_train', '''Online Learning
-----------------
• Streaming data
• Incremental updates
• Dynamic model''')
            
            sgd.node('sgd_benefit', '''Benefits
-----------------
• Low memory usage
• Quick updates
• Adapts to changes''')
        
        # Add comparison arrow
        dot.node('vs', 'VS', shape='circle', fillcolor='lightyellow')
        
        # Connect components
        dot.edge('rf_train', 'rf_limit')
        dot.edge('sgd_train', 'sgd_benefit')
    
    # Add simple note about concept drift
    dot.node('note1', '''Handling Concept Drift
-----------------
Random Forest: Requires retraining
SGD: Continuous adaptation''',
    shape='note', fillcolor='lightyellow')
    
    # Save diagram
    dot.attr(dpi='300')
    dot.attr(size='10,4')  # Smaller size for simpler diagram
    dot.render('presentation_images/sgd_diagram', format='png', cleanup=True)
    return Path('presentation_images/sgd_diagram.png')

def create_rf_methodology_diagram():
    """Create a basic diagram showing RandomForestClassifier methodology"""
    dot = graphviz.Digraph(comment='Random Forest Methodology')
    dot.attr(rankdir='TB')  # Top to Bottom direction
    dot.attr('node', shape='rectangle', style='rounded,filled')
    
    # Main process
    with dot.subgraph(name='cluster_0') as c:
        c.attr(label='Random Forest Training Methodology', style='rounded', color='navy', fontcolor='navy')
        
        # Input Data
        with dot.subgraph(name='cluster_input') as inp:
            inp.attr(label='Input', style='rounded', color='blue', fontcolor='blue')
            inp.attr('node', fillcolor='lightblue')
            
            inp.node('data', '''GitHub Issues
-----------------
Text Data''')
            
            inp.node('features', '''TF-IDF Features
-----------------
5000 dimensions''')
        
        # Random Forest Process
        with dot.subgraph(name='cluster_rf') as rf:
            rf.attr(label='Random Forest', style='rounded', color='darkgreen', fontcolor='darkgreen')
            rf.attr('node', fillcolor='lightgreen')
            
            # Create nodes for trees
            for i in range(3):  # Show 3 trees to represent 100
                rf.node(f'tree{i}', f'''Decision Tree {i+1}
-----------------
• Random Features
• Random Samples''')
            
            rf.node('ensemble', '''Ensemble
-----------------
100 Trees Total''')
            
            # Add dots to show more trees
            rf.node('dots', '...', shape='none')
        
        # Prediction Process
        with dot.subgraph(name='cluster_pred') as pred:
            pred.attr(label='Classification', style='rounded', color='darkred', fontcolor='darkred')
            pred.attr('node', fillcolor='lightpink')
            
            pred.node('voting', '''Majority Voting
-----------------
Each tree votes for
the issue type''')
            
            pred.node('output', '''Final Prediction
-----------------
• Bug
• Enhancement
• Question''')
        
        # Connect components
        dot.edge('data', 'features')
        for i in range(3):
            dot.edge('features', f'tree{i}')
            dot.edge(f'tree{i}', 'voting')
        dot.edge('features', 'dots')
        dot.edge('dots', 'voting')
        dot.edge('voting', 'output')
    
    # Add methodology note
    dot.node('note1', '''scikit-learn Implementation
-----------------
from sklearn.ensemble import 
RandomForestClassifier

rf = RandomForestClassifier(
    n_estimators=100,
    random_state=42
)''',
    shape='note', fillcolor='lightyellow')
    
    # Save diagram
    dot.attr(dpi='300')
    dot.attr(size='8,10')  # Vertical format
    dot.render('presentation_images/rf_methodology', format='png', cleanup=True)
    return Path('presentation_images/rf_methodology.png')

def create_custom_diagrams():
    """Create all custom diagrams"""
    images = {}
    
    # Create images directory
    images_dir = Path('presentation_images')
    images_dir.mkdir(exist_ok=True)
    
    # Create and save custom diagrams
    images['pipeline.png'] = create_pipeline_diagram()
    images['training_flow.png'] = create_training_flow_diagram()
    images['concept_drift.png'] = create_concept_drift_diagram()
    images['performance.png'] = create_performance_visualization()
    images['evaluation.png'] = create_evaluation_diagram()
    images['sgd.png'] = create_sgd_diagram()
    images['rf_methodology.png'] = create_rf_methodology_diagram()
    
    return images

def create_presentation(images):
    prs = Presentation()
    
    # Title Slide
    add_title_slide(
        prs,
        "Random Forest Model Pipeline for GitHub Issue Classification",
        "Exercise 3 - Principles of AI Engineering\nSyed Muhammed Hassan Ali"
    )
    
    # Pipeline Components Slide
    add_content_slide_with_image(
        prs,
        "Model Pipeline Components",
        [
            "Text Processing Pipeline:",
            "• URL & Code Block Removal",
            "• Text Normalization",
            "• Tokenization & Lemmatization",
            "Feature Extraction:",
            "• TF-IDF Vectorization",
            "Model Training:",
            "• Random Forest with Cross-Validation"
        ],
        images.get('pipeline.png')
    )
    
    # Training Flow Slide
    add_content_slide_with_image(
        prs,
        "Training Process Flow",
        [
            "1. Data Preparation",
            "• Load and clean GitHub issues",
            "• Apply text preprocessing steps",
            "2. Feature Engineering",
            "• Convert text to TF-IDF features",
            "3. Model Training",
            "• Train Random Forest model",
            "• Validate performance"
        ],
        images.get('training_flow.png')
    )
    
    # Training and Evaluation Slide
    add_content_slide_with_image(
        prs,
        "Training and Evaluation of the Model",
        [
            "Training Methodology:",
            "• Used scikit-learn's RandomForestClassifier",
            "• Applied cross-validation for reliable assessment",
            "• Serialized model using joblib",
            "",
            "Evaluation Metrics:",
            "• Overall Accuracy: 73%",
            "• Bug Classification:",
            "  - Precision: 76%, Recall: 79%",
            "• Enhancement Classification:",
            "  - Precision: 70%, Recall: 80%",
            "• Question Classification:",
            "  - Precision: 61%, Recall: 9%"
        ],
        images.get('evaluation.png')
    )
    
    # Model Performance Slide
    add_content_slide_with_image(
        prs,
        "Model Performance Analysis",
        [
            "Performance by Issue Type:",
            "• Strong performance on Bugs",
            "• Good results for Enhancements",
            "• Challenges with Questions",
            "Key Insights:",
            "• Class imbalance impact",
            "• Need for balanced dataset"
        ],
        images.get('performance.png')
    )
    
    # Concept Drift Slide
    add_content_slide_with_image(
        prs,
        "Impact of Concept Drift",
        [
            "Observed Changes:",
            "• Shifting issue patterns",
            "• Evolution of terminology",
            "• New feature requests",
            "Challenges:",
            "• Static model limitations",
            "• Need for adaptive learning"
        ],
        images.get('concept_drift.png')
    )
    
    # Add SGD Classifier Slide
    add_content_slide_with_image(
        prs,
        "Alternative: SGD Classifier for Concept Drift",
        [
            "Advantages:",
            "• Online Learning Capability",
            "  - Continuous model updates",
            "  - Adapts to new patterns",
            "• Memory Efficient",
            "  - No stored trees",
            "  - Linear model complexity",
            "• Scalable Solution",
            "  - Fast training and updates",
            "  - Handles streaming data"
        ],
        images.get('sgd.png')
    )
    
    # Save the presentation
    prs.save('random_forest_presentation_with_images.pptx')
    print("Presentation created successfully!")

def main():
    # Create custom diagrams
    print("Creating custom diagrams...")
    images = create_custom_diagrams()
    
    # Create presentation with images
    print("Creating presentation...")
    create_presentation(images)

if __name__ == "__main__":
    main() 